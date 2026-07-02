import os
from typing import List, Dict, Any
from pathlib import Path
from langchain_core.documents import Document
from utils.logger import logger
from .base_processor import BaseProcessor


class PPTXProcessor(BaseProcessor):
    """Process PPTX files and extract text from slides"""

    def process(self, file_path: str, metadata: Dict[str, Any] = None) -> List[Document]:
        """
        Process PPTX file
        
        Args:
            file_path: Path to PPTX file
            metadata: Optional metadata dictionary
            
        Returns:
            List of Document objects
        """
        documents = []
        metadata = metadata or {}
        
        try:
            from pptx import Presentation
            
            logger.info(f"Processing PPTX: {file_path}")
            
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    content = "\n".join(slide_text)
                    doc_metadata = {
                        **metadata,
                        "slide_number": slide_num,
                        "source": Path(file_path).name,
                    }
                    
                    doc = Document(
                        page_content=content,
                        metadata=doc_metadata
                    )
                    documents.append(doc)
            
            logger.info(f"Extracted {len(documents)} slides from PPTX")
            return documents
            
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {str(e)}")
            raise
